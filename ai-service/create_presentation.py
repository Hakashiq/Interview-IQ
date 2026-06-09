import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_deck():
    prs = Presentation()
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Theme colors (Harmony Dark Palette matching the app UI)
    bg_color = RGBColor(15, 23, 42)      # Deep Slate Blue (#0f172a)
    title_color = RGBColor(14, 165, 233)  # Electric Cyan (#0ea5e9)
    accent_color = RGBColor(168, 85, 247) # Vibrant Purple (#a855f7)
    text_color = RGBColor(255, 255, 255)  # White (#ffffff)
    muted_color = RGBColor(148, 163, 184) # Light Gray (#94a3b8)
    
    blank_layout = prs.slide_layouts[6]
    
    slides_data = [
        # Slide 1: Title
        {
            "type": "title",
            "title": "InterviewIQ",
            "subtitle": "Enterprise-Grade AI Mock Interview & Resume Analyzer Platform",
            "info": "Platform Features, Architecture, Security Framework, and Core Capabilities\nCreated by: Google DeepMind team pair programming with USER"
        },
        # Slide 2: Executive Summary
        {
            "type": "content",
            "title": "Executive Summary",
            "subtitle": "Closing the Tech Skill-Gap with Intelligent Automation",
            "points": [
                "Advanced AI-driven mock interview environment offering realistic, domain-specific scenarios.",
                "Automated resume parsing and ATS matching rule-engine validation for target job roles.",
                "Detailed, structured evaluation rubrics (technical, communication, problem-solving, behavioral).",
                "Full-featured platform dashboard, student activity tracking, and comprehensive administrator configuration hub.",
                "Designed with strict security monitoring, fraud logs compliance, and secure Role-Based Access Control (RBAC)."
            ]
        },
        # Slide 3: Tech Stack & Architecture
        {
            "type": "content",
            "title": "Enterprise System Architecture",
            "subtitle": "Scalable Frameworks and Real-Time Service Topologies",
            "points": [
                "Frontend: Responsive React SPA, styled using custom modern dark-theme CSS variables, leveraging Chart.js/react-chartjs-2 for metrics and stats visualization.",
                "Backend: Enterprise-grade Spring Boot 3 web services. Manages auth (JWT), database operations (JPA + Hibernate), and business logic.",
                "Database Layer: MySQL Server storing persistent data (Users, Roles, Questions, SystemConfig, Audit Logs, Violations, Feedback).",
                "AI Service Layer: Python FastAPI serving as client bridge to Google Gemini Pro API for prompt engineering and evaluation models.",
                "Integration: Asynchronous REST communications between subsystems with robust fallback heuristics."
            ]
        },
        # Slide 4: AI Mock Interview Engine
        {
            "type": "content",
            "title": "AI Mock Interview Engine",
            "subtitle": "Interactive, Dynamic Session Generation",
            "points": [
                "Adaptive Difficulty: Shifts questions from Easy to Hard dynamically based on answer scores of preceding rounds.",
                "Multiple Assessment Modes: Supports text-input answering with upcoming expansion for voice-to-text response synthesis.",
                "Role-Based Generation: Predefined system templates for Backend Developer, Full Stack, Software Engineer, DevOps, and Data Analyst.",
                "Question Bank: Comprehensive domain bank including Java, Python, DSA, DBMS, System Design, React, and behavioral categories.",
                "Dual Personality AI: Configurable examiner personalities (Professional, Strict, Friendly, Challenging) to test pressure composure."
            ]
        },
        # Slide 5: Multi-Dimensional AI Evaluations
        {
            "type": "content",
            "title": "Multi-Dimensional AI Evaluation",
            "subtitle": "Comprehensive Feedback & Score Remarks",
            "points": [
                "Technical Accuracy: Checks correctness of technical facts, syntax, and conceptual foundations.",
                "Completeness & Detail: Audits depth of explanations, missing edge-cases, and structural completeness.",
                "Communication & Vocabulary: Scans fluency, readability, filler word checks, and clarity of explanation.",
                "Problem Solving: Evaluates design choices, algorithmic complexity (Big O), and structural breakdown.",
                "Actionable Remarks: Returns precise highlighted mistakes, strengths, weaknesses, and concrete recommendations for improvement."
            ]
        },
        # Slide 6: ATS Resume Scanner
        {
            "type": "content",
            "title": "ATS Rule Engine & Resume Analyzer",
            "subtitle": "Automating Match Verification against Target Job Profiles",
            "points": [
                "Multi-Format Parser: Uses Python library parsing engines (PyMuPDF, docx) to extract text content dynamically.",
                "Keyword Matching: Scores resume text keyword matching frequencies against role-specific requirements.",
                "Format & Structure Audits: Scans for parsing-safe layout rules, contact information, education, and experience details.",
                "Dynamic Weighting: Configurable scoring weights (Keywords, Projects, Skill Validation, Experience, Education) adjusted in Admin panel.",
                "ATS passing threshold: Real-time validation scoring matching candidate resumes to target roles (e.g. Backend vs DevOps)."
            ]
        },
        # Slide 7: Centralized Admin Settings Center
        {
            "type": "content",
            "title": "Centralized Administration & Configurations",
            "subtitle": "Eliminating Hardcoded Defaults with Dynamic MySQL Schemas",
            "points": [
                "Dynamic Configuration: Admin settings read and write to key-value `system_config` table (platform name, logo, models, etc.).",
                "Prompt Manager: Admin controls and versions prompts for Interview Generation, Evaluations, ATS, and Recommendations with rollbacks.",
                "Interview Template Builder: Dynamically configure question counts, durations, and skills requirements for new roles.",
                "Audit Logging: Automatically tracks every administrative action, user modification, ban, and system update.",
                "Resource Management: Monitor token usage metrics, rate limits, daily quotas, and estimated API provider costs."
            ]
        },
        # Slide 8: Security, Compliance & Fraud Prevention
        {
            "type": "content",
            "title": "Security, Compliance & Fraud Monitoring",
            "subtitle": "Protecting Platform Integrity and Session Authenticity",
            "points": [
                "Access Control: Secure JWT token authentication. Granular roles: Student, Admin, Super Admin, Mentor, Analyst, Moderator.",
                "Two-Factor Authentication: Toggleable 2FA security enforcement for administrator accounts.",
                "Copy-Paste Prevention: Listeners detect copy-paste events during active interview sessions and log violations.",
                "Tab-Switch Monitoring: Tracking active window shifts during sessions, auto-flagging suspicious activities.",
                "Integrity Scores: Calculates system integrity metrics and logs violation reports to database audit tables."
            ]
        },
        # Slide 9: Collaborative Mentor Workspace
        {
            "type": "content",
            "title": "Collaborative Mentor Hub",
            "subtitle": "Empowering Educators and Advisors to Track Student Success",
            "points": [
                "Student Directory: Comprehensive search/filter list of student accounts, education info, addresses, and resumes.",
                "Performance Tracking: Monitor student interview success rates, average scores, and weak preparation topics.",
                "Complaints & Feedback: Reviews student suggestions and complaints directly with toggle indicators.",
                "Integrity Audits: Mentors review integrity infraction history (copy-paste, tab switches) for assigned students.",
                "Personalized Roadmap Reviews: Mentor-guided prep recommendations for students based on AI insights."
            ]
        },
        # Slide 10: Conclusion & Strategic Roadmap
        {
            "type": "content",
            "title": "Conclusion & Strategic Roadmap",
            "subtitle": "Future Extensions and Platform Scope",
            "points": [
                "Fully operational MySQL database integration replacing all mocked/hardcoded configurations.",
                "Completed Sidebar layout with collapse/expand triggers and customized tooltips.",
                "LeetCode-style daily streak flames tracker on student dashboard to drive daily practice engagement.",
                "Future Scope: Voice Synthesis (TTS) and voice response analysis for video-interview simulations.",
                "Future Scope: Real-time code execution compiler and multi-candidate online whiteboard coding rounds."
            ]
        }
    ]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        if slide_data["type"] == "title":
            # Add logo block
            tx_logo = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1))
            tf_logo = tx_logo.text_frame
            tf_logo.word_wrap = True
            p_logo = tf_logo.paragraphs[0]
            p_logo.text = "INTERVIEW IQ"
            p_logo.font.name = "Segoe UI"
            p_logo.font.size = Pt(20)
            p_logo.font.bold = True
            p_logo.font.color.rgb = title_color
            
            # Title
            tx_title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
            tf_title = tx_title.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = slide_data["title"]
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(64)
            p_title.font.bold = True
            p_title.font.color.rgb = text_color
            
            # Subtitle
            tx_sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
            tf_sub = tx_sub.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = slide_data["subtitle"]
            p_sub.font.name = "Segoe UI"
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = accent_color
            
            # Muted Info
            tx_info = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1.5))
            tf_info = tx_info.text_frame
            tf_info.word_wrap = True
            p_info = tf_info.paragraphs[0]
            p_info.text = slide_data["info"]
            p_info.font.name = "Segoe UI"
            p_info.font.size = Pt(14)
            p_info.font.color.rgb = muted_color
            
        else:
            # Header section
            tx_header = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.5))
            tf_header = tx_header.text_frame
            tf_header.word_wrap = True
            tf_header.margin_top = Inches(0.1)
            tf_header.margin_bottom = Inches(0.1)
            
            # Category / Section title
            p_sec = tf_header.paragraphs[0]
            p_sec.text = "INTERVIEW IQ  |  " + slide_data["title"].upper()
            p_sec.font.name = "Segoe UI"
            p_sec.font.size = Pt(12)
            p_sec.font.bold = True
            p_sec.font.color.rgb = title_color
            p_sec.space_after = Pt(8)
            
            # Main Title
            p_title = tf_header.add_paragraph()
            p_title.text = slide_data["title"]
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(36)
            p_title.font.bold = True
            p_title.font.color.rgb = text_color
            
            # Subtitle
            p_sub = tf_header.add_paragraph()
            p_sub.text = slide_data["subtitle"]
            p_sub.font.name = "Segoe UI"
            p_sub.font.size = Pt(16)
            p_sub.font.color.rgb = accent_color
            
            # Content Bullets
            tx_content = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(4.5))
            tf_content = tx_content.text_frame
            tf_content.word_wrap = True
            
            for i, pt in enumerate(slide_data["points"]):
                p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
                p.text = "•  " + pt
                p.font.name = "Segoe UI"
                p.font.size = Pt(16)
                p.font.color.rgb = text_color
                p.space_after = Pt(14)
                p.line_spacing = 1.25
                
    # Save the presentation in the interview-iq base directory
    output_path = r"d:\Antigravity\interview-iq\presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == '__main__':
    create_deck()
